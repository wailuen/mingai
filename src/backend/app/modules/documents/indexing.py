"""
DocumentIndexingPipeline (AI-060).

Parses documents (PDF, DOCX, PPTX, TXT), chunks text into ~512-token segments
with 50-token overlap, embeds each chunk, and upserts to the vector store.

Supported extensions: .pdf, .docx, .pptx, .txt
"""
import asyncio
import os

import structlog
from sqlalchemy import text
from sqlalchemy.ext.asyncio import AsyncSession

from app.modules.chat.embedding import EmbeddingService
from app.modules.chat.vector_search import VectorSearchService

logger = structlog.get_logger()


class DocumentIndexingPipeline:
    """
    Parse → chunk → embed → vector upsert pipeline for documents.

    Chunk size: 512 tokens (approximated as 2048 chars, ~4 chars/token)
    Chunk overlap: 50 tokens (200 chars)
    """

    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}
    CHUNK_SIZE = 512  # tokens
    CHUNK_OVERLAP = 50  # tokens

    # Character approximations (4 chars ≈ 1 token)
    _CHUNK_SIZE_CHARS = CHUNK_SIZE * 4  # 2048
    _CHUNK_OVERLAP_CHARS = CHUNK_OVERLAP * 4  # 200

    # Zip bomb guard: cap extracted text at 1M chars
    _MAX_TEXT_CHARS = 1_000_000

    async def process_file(
        self,
        file_path: str,
        integration_id: str,
        tenant_id: str,
        db: AsyncSession,
    ) -> dict:
        """
        Index a single file:
          1. Detect file type by extension
          2. Extract text
          3. Chunk text
          4. Embed each chunk
          5. Upsert chunks to vector store
          6. Update sync_jobs record
          7. Return result dict

        Raises ValueError for unsupported extensions.
        """
        ext = os.path.splitext(file_path)[1].lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file extension '{ext}'. "
                f"Supported: {sorted(self.SUPPORTED_EXTENSIONS)}"
            )

        # Extract text
        if ext == ".pdf":
            text_content = self._extract_text_pdf(file_path)
        elif ext == ".docx":
            text_content = self._extract_text_docx(file_path)
        elif ext == ".pptx":
            text_content = self._extract_text_pptx(file_path)
        else:  # .txt
            text_content = self._extract_text_txt(file_path)

        # Chunk text
        chunks = self._chunk_text(text_content)
        chunks_indexed = 0

        if chunks:
            embedding_service = EmbeddingService()
            vector_service = VectorSearchService()

            for i, chunk in enumerate(chunks):
                try:
                    embedding = await embedding_service.embed(
                        chunk, tenant_id=tenant_id
                    )
                    await vector_service.upsert_chunks(
                        tenant_id=tenant_id,
                        integration_id=integration_id,
                        file_path=file_path,
                        chunk_index=i,
                        chunk_text=chunk,
                        embedding=embedding,
                    )
                    chunks_indexed += 1
                except Exception as exc:
                    logger.error(
                        "document_indexing_chunk_error",
                        file_path=file_path,
                        chunk_index=i,
                        error=str(exc),
                    )

        # Update sync_jobs: increment files_synced or files_failed
        if chunks_indexed > 0:
            await db.execute(
                text(
                    "UPDATE sync_jobs "
                    "SET files_synced = files_synced + 1 "
                    "WHERE integration_id = :integration_id "
                    "AND tenant_id = :tenant_id "
                    "AND status IN ('running', 'pending')"
                ),
                {"integration_id": integration_id, "tenant_id": tenant_id},
            )
        else:
            await db.execute(
                text(
                    "UPDATE sync_jobs "
                    "SET files_failed = files_failed + 1 "
                    "WHERE integration_id = :integration_id "
                    "AND tenant_id = :tenant_id "
                    "AND status IN ('running', 'pending')"
                ),
                {"integration_id": integration_id, "tenant_id": tenant_id},
            )

        logger.info(
            "document_indexed",
            file_path=file_path,
            integration_id=integration_id,
            tenant_id=tenant_id,
            chunks_indexed=chunks_indexed,
            total_chunks=len(chunks),
        )

        return {
            "file_path": file_path,
            "chunks_indexed": chunks_indexed,
            "integration_id": integration_id,
        }

    async def process_conversation_file(
        self,
        file_path: str,
        conversation_id: str,
        user_id: str,
        tenant_id: str,
        original_filename: str | None = None,
    ) -> dict:
        """
        Index a user-uploaded conversation document.

        Unlike process_file(), this method:
        - Uses conversation-scoped index_id (conv-{tenant_id}-{conv_id})
        - Does NOT write to sync_jobs (no integration involved)
        - Caps extracted text at 1M chars (zip bomb guard)
        - Offloads synchronous text extraction to asyncio.to_thread()

        Returns:
            {"chunks_indexed": N, "file_name": basename, "conversation_id": ..., "index_id": ...}
        """
        ext = os.path.splitext(file_path)[1].lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file extension '{ext}'. "
                f"Supported: {sorted(self.SUPPORTED_EXTENSIONS)}"
            )

        # Extract text in thread pool to avoid blocking event loop
        if ext == ".pdf":
            text_content = await asyncio.to_thread(self._extract_text_pdf, file_path)
        elif ext == ".docx":
            text_content = await asyncio.to_thread(self._extract_text_docx, file_path)
        elif ext == ".pptx":
            text_content = await asyncio.to_thread(self._extract_text_pptx, file_path)
        else:  # .txt
            text_content = await asyncio.to_thread(self._extract_text_txt, file_path)

        # Cap extracted text (zip bomb guard)
        if len(text_content) > self._MAX_TEXT_CHARS:
            logger.warning(
                "conversation_file_text_capped",
                file_path=file_path,
                original_length=len(text_content),
                capped_length=self._MAX_TEXT_CHARS,
            )
            text_content = text_content[: self._MAX_TEXT_CHARS]

        # Chunk text
        chunks = self._chunk_text(text_content)

        if not chunks:
            return {
                "chunks_indexed": 0,
                "file_name": original_filename or os.path.basename(file_path),
                "conversation_id": conversation_id,
                "index_id": f"conv-{tenant_id}-{conversation_id}",
            }

        embedding_service = EmbeddingService()
        vector_service = VectorSearchService()
        # Use caller-supplied original filename if provided (avoids exposing temp path)
        file_name = original_filename or os.path.basename(file_path)

        # Build all embeddings and chunk dicts, then upsert in one batch
        chunk_dicts = []
        chunks_indexed = 0
        for i, chunk in enumerate(chunks):
            try:
                embedding = await embedding_service.embed(chunk, tenant_id=tenant_id)
                chunk_dicts.append({"text": chunk, "embedding": embedding})
            except Exception as exc:
                logger.error(
                    "conversation_doc_chunk_embed_error",
                    file_path=file_path,
                    chunk_index=i,
                    error=str(exc),
                )

        if chunk_dicts:
            try:
                await vector_service.upsert_conversation_chunks(
                    tenant_id=tenant_id,
                    conversation_id=conversation_id,
                    user_id=user_id,
                    file_name=file_name,
                    chunks=chunk_dicts,
                )
                chunks_indexed = len(chunk_dicts)
            except Exception as exc:
                logger.error(
                    "conversation_doc_upsert_error",
                    file_path=file_path,
                    error=str(exc),
                )

        logger.info(
            "conversation_document_indexed",
            file_path=file_path,
            conversation_id=conversation_id,
            tenant_id=tenant_id,
            chunks_indexed=chunks_indexed,
            total_chunks=len(chunks),
        )

        return {
            "chunks_indexed": chunks_indexed,
            "file_name": file_name,
            "conversation_id": conversation_id,
            "index_id": f"conv-{tenant_id}-{conversation_id}",
        }

    def _extract_text_pdf(self, file_path: str) -> str:
        """Extract text from PDF using pypdf."""
        try:
            from pypdf import PdfReader
        except ImportError:
            raise ImportError(
                "pypdf is required for PDF extraction. "
                "Add 'pypdf>=4.0.0' to project dependencies."
            )
        reader = PdfReader(file_path)
        parts = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                parts.append(page_text)
        return "\n".join(parts)

    def _extract_text_docx(self, file_path: str) -> str:
        """Extract text from DOCX using python-docx."""
        try:
            from docx import Document
        except ImportError:
            raise ImportError(
                "python-docx is required for DOCX extraction. "
                "Add 'python-docx>=1.1.0' to project dependencies."
            )
        doc = Document(file_path)
        return "\n".join(para.text for para in doc.paragraphs if para.text.strip())

    def _extract_text_pptx(self, file_path: str) -> str:
        """Extract text from PPTX using python-pptx. Extracts text from all slides."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required for PPTX extraction. "
                "Add 'python-pptx>=0.6.23' to project dependencies."
            )
        prs = Presentation(file_path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n".join(parts)

    def _extract_text_txt(self, file_path: str) -> str:
        """Read plain text file, detect encoding with chardet."""
        try:
            import chardet
        except ImportError:
            # Fallback to utf-8 with errors='replace'
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()

        with open(file_path, "rb") as f:
            raw = f.read()

        detected = chardet.detect(raw)
        encoding = detected.get("encoding") or "utf-8"
        return raw.decode(encoding, errors="replace")

    def _chunk_text(self, text: str) -> list[str]:
        """
        Split text into chunks of CHUNK_SIZE*4 chars with CHUNK_OVERLAP*4 char overlap.

        Returns a list of string chunks. Single chunk if text fits.
        """
        chunk_size = self._CHUNK_SIZE_CHARS
        overlap = self._CHUNK_OVERLAP_CHARS

        if len(text) <= chunk_size:
            return [text] if text.strip() else []

        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            if chunk.strip():
                chunks.append(chunk)
            if end >= len(text):
                break
            start = end - overlap

        return chunks
